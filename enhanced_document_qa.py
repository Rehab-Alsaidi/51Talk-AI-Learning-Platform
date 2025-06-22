import os
import logging
import threading
import time
from typing import List, Dict, Any, Optional
import pptx
import re
from datetime import datetime

from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

# Add Llama integration
try:
    from llama_cpp import Llama
    LLAMA_AVAILABLE = True
except ImportError:
    LLAMA_AVAILABLE = False
    print("Warning: llama-cpp-python not installed. Install with: pip install llama-cpp-python")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class LlamaLLM:
    """Llama model wrapper for local inference."""
    
    def __init__(self, model_path: str, **kwargs):
        if not LLAMA_AVAILABLE:
            raise ImportError("llama-cpp-python is required. Install with: pip install llama-cpp-python")
        
        self.model_path = model_path
        self.llm = None
        self._initialize_model(**kwargs)
    
    def _initialize_model(self, **kwargs):
        """Initialize the Llama model with optimized settings."""
        default_kwargs = {
            'n_ctx': 4096,  # Context window
            'n_threads': 8,  # Number of threads
            'n_gpu_layers': 0,  # Set to -1 for full GPU offload if you have CUDA
            'verbose': False,
            'temperature': 0.1,  # Low temperature for more focused responses
            'max_tokens': 512,  # Maximum response length
            'top_p': 0.9,
            'repeat_penalty': 1.1,
        }
        
        # Merge user kwargs with defaults
        model_kwargs = {**default_kwargs, **kwargs}
        
        try:
            logger.info(f"Loading Llama model from {self.model_path}")
            self.llm = Llama(model_path=self.model_path, **model_kwargs)
            logger.info("Llama model loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load Llama model: {str(e)}")
            raise
    
    def generate_response(self, prompt: str, **kwargs) -> str:
        """Generate response using the Llama model."""
        if not self.llm:
            raise RuntimeError("Llama model not initialized")
        
        try:
            response = self.llm(prompt, **kwargs)
            return response['choices'][0]['text'].strip()
        except Exception as e:
            logger.error(f"Error generating response: {str(e)}")
            return "I apologize, but I encountered an error while generating a response."

class PPTXLoader:
    """Enhanced PPTX loader with better error handling."""
    
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Dict[str, Any]]:
        """Load PPTX documents with improved text extraction."""
        docs = []
        try:
            prs = pptx.Presentation(self.file_path)
            filename = os.path.basename(self.file_path)
            
            for i, slide in enumerate(prs.slides):
                text_content = ""
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    
                    # Handle tables if present
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_content += cell.text + " "
                            text_content += "\n"
                
                if text_content.strip():
                    metadata = {
                        "source": self.file_path,
                        "slide_number": i + 1,
                        "page": i + 1,  # For consistency with PDF
                        "filename": filename
                    }
                    doc = {
                        "page_content": text_content.strip(),
                        "metadata": metadata
                    }
                    docs.append(doc)
                    
            logger.info(f"Successfully loaded {len(docs)} slides from {filename}")
            return docs
            
        except Exception as e:
            logger.error(f"Error processing PPTX {self.file_path}: {str(e)}")
            return []

class DocumentProcessor:
    """Enhanced document processor with async capabilities."""
    
    def __init__(self, documents_dir: str):
        self.documents_dir = documents_dir
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,  # Reduced for better performance
            chunk_overlap=100,
            length_function=len,
            is_separator_regex=False
        )
        self._loading_complete = False
        self._documents_cache = None

    def load_documents(self) -> List[Any]:
        """Load all documents with improved error handling."""
        if self._documents_cache is not None and self._loading_complete:
            return self._documents_cache
            
        all_docs = []
        supported_extensions = {'.pdf', '.pptx', '.ppt', '.txt'}
        
        if not os.path.exists(self.documents_dir):
            logger.warning(f"Documents directory not found: {self.documents_dir}")
            return []
        
        total_files = 0
        processed_files = 0
        
        # Count total files first
        for root, _, files in os.walk(self.documents_dir):
            for file in files:
                if any(file.lower().endswith(ext) for ext in supported_extensions):
                    total_files += 1
        
        logger.info(f"Found {total_files} supported documents to process")
        
        for root, _, files in os.walk(self.documents_dir):
            for file in files:
                file_path = os.path.join(root, file)
                filename = os.path.basename(file_path)
                
                try:
                    if file.lower().endswith('.pdf'):
                        docs = self._load_pdf(file_path, filename)
                        all_docs.extend(docs)
                        
                    elif file.lower().endswith(('.pptx', '.ppt')):
                        docs = self._load_pptx(file_path)
                        all_docs.extend(docs)
                        
                    elif file.lower().endswith('.txt'):
                        docs = self._load_txt(file_path, filename)
                        all_docs.extend(docs)
                    
                    processed_files += 1
                    if processed_files % 5 == 0:  # Log progress every 5 files
                        logger.info(f"Processed {processed_files}/{total_files} files...")
                        
                except Exception as e:
                    logger.error(f"Error processing {file_path}: {str(e)}")
                    continue
        
        logger.info(f"Successfully loaded {len(all_docs)} document sections from {processed_files} files")
        self._documents_cache = all_docs
        self._loading_complete = True
        return all_docs

    def _load_pdf(self, file_path: str, filename: str) -> List[Any]:
        """Load PDF with enhanced metadata."""
        try:
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            
            # Enhance metadata
            for doc in docs:
                doc.metadata['filename'] = filename
                
            logger.debug(f"Loaded PDF: {filename} ({len(docs)} pages)")
            return docs
            
        except Exception as e:
            logger.error(f"Error loading PDF {filename}: {str(e)}")
            return []

    def _load_pptx(self, file_path: str) -> List[Any]:
        """Load PPTX with enhanced processing."""
        try:
            loader = PPTXLoader(file_path)
            return loader.load()
        except Exception as e:
            logger.error(f"Error loading PPTX {file_path}: {str(e)}")
            return []

    def _load_txt(self, file_path: str, filename: str) -> List[Any]:
        """Load TXT files with encoding detection."""
        encodings = ['utf-8', 'utf-16', 'iso-8859-1', 'windows-1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                if content.strip():
                    doc = {
                        "page_content": content.strip(),
                        "metadata": {
                            "source": file_path,
                            "page": 1,
                            "filename": filename
                        }
                    }
                    logger.debug(f"Loaded TXT: {filename}")
                    return [doc]
                    
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error loading TXT {filename}: {str(e)}")
                break
                
        return []

    def process_documents(self) -> List[Any]:
        """Process documents into chunks with progress tracking."""
        docs = self.load_documents()
        
        if not docs:
            logger.warning("No documents were loaded. Check your directory path and file formats.")
            return []
        
        logger.info(f"Processing {len(docs)} document sections into chunks...")
        
        try:
            # Extract content and metadata
            texts = []
            metadatas = []
            
            for doc in docs:
                if hasattr(doc, "page_content"):
                    content = doc.page_content
                    metadata = doc.metadata
                else:
                    content = doc["page_content"]
                    metadata = doc["metadata"]
                
                # Clean content
                content = self._clean_content(content)
                if len(content.strip()) > 50:  # Only include substantial content
                    texts.append(content)
                    metadatas.append(metadata)
            
            # Create chunks
            splits = self.text_splitter.create_documents(texts, metadatas=metadatas)
            logger.info(f"Created {len(splits)} document chunks from {len(texts)} sections")
            
            return splits
            
        except Exception as e:
            logger.error(f"Error processing documents: {str(e)}")
            return []

    def _clean_content(self, content: str) -> str:
        """Clean and normalize document content."""
        # Remove excessive whitespace
        content = re.sub(r'\s+', ' ', content)
        
        # Remove page numbers and headers/footers
        lines = content.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            # Skip likely page numbers, headers, footers
            if (len(line) < 5 or 
                re.match(r'^\d+$', line) or 
                re.match(r'^Page \d+', line, re.IGNORECASE)):
                continue
            cleaned_lines.append(line)
        
        return ' '.join(cleaned_lines).strip()

class VectorStore:
    """Enhanced vector store with async initialization."""
    
    def __init__(self, documents_dir: str, vector_db_path: str = "vector_db"):
        self.documents_dir = documents_dir
        self.vector_db_path = vector_db_path
        self.embeddings = None
        self.vector_store_instance = None
        self._initialization_lock = threading.Lock()
        self._is_initializing = False
        self._initialization_complete = False
        self._initialization_error = None
        
        # Start initialization in background
        self._start_background_initialization()

    def _start_background_initialization(self):
        """Start background initialization of embeddings and vector store."""
        def initialize():
            try:
                with self._initialization_lock:
                    if self._initialization_complete or self._is_initializing:
                        return
                    
                    self._is_initializing = True
                    logger.info("Starting background initialization of vector store...")
                    
                    # Initialize embeddings
                    self.embeddings = HuggingFaceEmbeddings(
                        model_name="all-MiniLM-L6-v2",
                        model_kwargs={'device': 'cpu'},
                        encode_kwargs={'normalize_embeddings': True}
                    )
                    logger.info("Embeddings initialized successfully")
                    
                    # Load or create vector store
                    self.vector_store_instance = self._load_or_create_vector_store()
                    
                    self._initialization_complete = True
                    self._is_initializing = False
                    logger.info("Vector store initialization completed successfully")
                    
            except Exception as e:
                self._initialization_error = str(e)
                self._is_initializing = False
                logger.error(f"Failed to initialize vector store: {str(e)}")
        
        # Run in background thread
        thread = threading.Thread(target=initialize, daemon=True)
        thread.start()

    def _load_or_create_vector_store(self):
        """Load existing vector store or create new one."""
        if os.path.exists(self.vector_db_path):
            try:
                logger.info(f"Loading existing vector database from {self.vector_db_path}")
                return FAISS.load_local(
                    self.vector_db_path, 
                    self.embeddings, 
                    allow_dangerous_deserialization=True
                )
            except Exception as e:
                logger.warning(f"Failed to load existing vector store: {str(e)}. Creating new one.")
        
        return self._create_new_vector_store()

    def _create_new_vector_store(self):
        """Create a new vector store from documents."""
        processor = DocumentProcessor(self.documents_dir)
        document_chunks = processor.process_documents()
        
        if not document_chunks:
            logger.warning("No document chunks available for vector store creation")
            return None
        
        logger.info(f"Creating new vector database from {len(document_chunks)} chunks...")
        
        try:
            vector_store = FAISS.from_documents(
                documents=document_chunks,
                embedding=self.embeddings
            )
            
            # Save for future use
            vector_store.save_local(self.vector_db_path)
            logger.info(f"Vector database created and saved to {self.vector_db_path}")
            
            return vector_store
            
        except Exception as e:
            logger.error(f"Failed to create vector database: {str(e)}")
            return None

    def get_vector_store(self, timeout: int = 30):
        """Get vector store with timeout for initialization."""
        start_time = time.time()
        
        while time.time() - start_time < timeout:
            if self._initialization_complete:
                return self.vector_store_instance
            elif self._initialization_error:
                raise Exception(f"Vector store initialization failed: {self._initialization_error}")
            
            time.sleep(0.5)  # Check every 500ms
        
        raise TimeoutError("Vector store initialization timed out")

    def is_ready(self) -> bool:
        """Check if vector store is ready for use."""
        return self._initialization_complete and self.vector_store_instance is not None

class DocumentQA:
    """Enhanced Document QA with Llama integration."""
    
    def __init__(self, vector_store=None, vector_store_manager=None, llama_model_path=None):
        self.vector_store = vector_store
        self.vector_store_manager = vector_store_manager
        self.llama_llm = None
        
        # Initialize Llama model if path is provided
        if llama_model_path and LLAMA_AVAILABLE:
            try:
                self.llama_llm = LlamaLLM(llama_model_path)
                logger.info("Llama model initialized successfully")
            except Exception as e:
                logger.error(f"Failed to initialize Llama model: {str(e)}")
        elif llama_model_path and not LLAMA_AVAILABLE:
            logger.warning("Llama model path provided but llama-cpp-python not available")
            
        logger.info("DocumentQA initialized with Llama integration")

    def _create_llama_prompt(self, question: str, context: str) -> str:
        """Create a structured prompt for the Llama model."""
        prompt = f"""<|begin_of_text|><|start_header_id|>system<|end_header_id|>

You are a helpful AI assistant that answers questions based on provided course materials. Use only the information from the context to answer questions. If the context doesn't contain enough information to answer the question, say so clearly.

Guidelines:
- Be concise but comprehensive
- Use bullet points or numbered lists when appropriate
- Cite specific details from the materials when possible
- If you cannot answer based on the context, suggest what additional information might be needed<|eot_id|>

<|start_header_id|>user<|end_header_id|>

Context from course materials:
{context}

Question: {question}<|eot_id|>

<|start_header_id|>assistant<|end_header_id|>

"""
        return prompt

    def _extract_sources_from_docs(self, docs) -> List[Dict[str, str]]:
        """Extract source information from documents."""
        sources = []
        seen_sources = set()
        
        for doc in docs:
            metadata = doc.metadata if hasattr(doc, "metadata") else doc.get("metadata", {})
            filename = metadata.get("filename", os.path.basename(metadata.get("source", "Unknown")))
            page = metadata.get("page", metadata.get("slide_number", "Unknown"))
            
            source_key = f"{filename}_{page}"
            if source_key not in seen_sources:
                sources.append({
                    "file": filename,
                    "page": str(page)
                })
                seen_sources.add(source_key)
        
        return sources[:4]  # Limit to 4 sources

    def answer_question(self, question: str) -> Dict[str, Any]:
        """Generate responses using Llama model with document context."""
        if not question or not question.strip():
            return {
                "answer": "Please provide a question for me to answer.",
                "sources": []
            }
        
        try:
            # Get vector store
            vector_store = None
            
            if self.vector_store_manager:
                if not self.vector_store_manager.is_ready():
                    if self.vector_store_manager._is_initializing:
                        return {
                            "answer": "I'm currently loading the course materials. Please wait a moment and try again.",
                            "sources": []
                        }
                    else:
                        return {
                            "answer": "The course materials are not available right now. Please contact your administrator.",
                            "sources": []
                        }
                
                vector_store = self.vector_store_manager.get_vector_store(timeout=10)
            else:
                vector_store = self.vector_store
            
            if not vector_store:
                return {
                    "answer": "I don't have access to course materials at the moment. Please ensure documents are uploaded and try again.",
                    "sources": []
                }
            
            # Retrieve relevant documents
            retriever = vector_store.as_retriever(
                search_type="similarity",
                search_kwargs={"k": 5}
            )
            
            docs = retriever.get_relevant_documents(question)
            
            if not docs:
                return {
                    "answer": "I couldn't find relevant information in the course materials to answer your question. Please try rephrasing your question or ask about topics covered in your materials.",
                    "sources": []
                }
            
            # Prepare context from retrieved documents
            context = ""
            for doc in docs:
                content = doc.page_content if hasattr(doc, "page_content") else doc.get("page_content", "")
                context += content + "\n\n"
            
            # Limit context length to avoid token limits
            if len(context) > 3000:
                context = context[:3000] + "..."
            
            sources = self._extract_sources_from_docs(docs)
            
            # Generate response using Llama if available, otherwise use rule-based approach
            if self.llama_llm:
                try:
                    prompt = self._create_llama_prompt(question, context)
                    answer = self.llama_llm.generate_response(
                        prompt,
                        max_tokens=512,
                        temperature=0.1,
                        top_p=0.9,
                        stop=["<|eot_id|>", "<|end_of_text|>"]
                    )
                    
                    # Clean up the response
                    answer = answer.strip()
                    if not answer:
                        answer = "I apologize, but I couldn't generate a proper response based on the available materials."
                    
                except Exception as e:
                    logger.error(f"Llama generation error: {str(e)}")
                    answer = "I encountered an error while generating a response. Please try again."
            else:
                # Fallback to rule-based response generation
                answer = self._generate_fallback_response(context, question)
            
            return {
                "answer": answer,
                "sources": sources
            }
            
        except Exception as e:
            logger.error(f"Error in answer_question: {str(e)}")
            return {
                "answer": "I encountered an unexpected error. Please try rephrasing your question or contact support if the problem persists.",
                "sources": []
            }

    def _generate_fallback_response(self, context: str, question: str) -> str:
        """Generate a simple response when Llama is not available."""
        # Extract the most relevant sentences
        sentences = [s.strip() for s in context.split('.') if s.strip() and len(s.strip()) > 20]
        
        # Find sentences that might contain the answer
        question_words = question.lower().split()
        relevant_sentences = []
        
        for sentence in sentences[:10]:  # Check first 10 sentences
            sentence_lower = sentence.lower()
            score = sum(1 for word in question_words if word in sentence_lower)
            if score > 0:
                relevant_sentences.append((sentence, score))
        
        # Sort by relevance and take top 3
        relevant_sentences.sort(key=lambda x: x[1], reverse=True)
        top_sentences = [s[0] for s in relevant_sentences[:3]]
        
        if top_sentences:
            return "Based on the course materials:\n\n" + ". ".join(top_sentences) + "."
        else:
            return "I found some relevant information in the course materials, but couldn't extract a specific answer to your question. Please try rephrasing your question or ask about more specific topics."

# Global instance management
_vector_store_manager = None
_document_qa = None

def initialize_document_qa(documents_dir: str, llama_model_path: str = None) -> DocumentQA:
    """Initialize DocumentQA with Llama integration."""
    global _vector_store_manager, _document_qa
    
    if _document_qa is None:
        logger.info(f"Initializing DocumentQA with documents from: {documents_dir}")
        if llama_model_path:
            logger.info(f"Using Llama model: {llama_model_path}")
        
        _vector_store_manager = VectorStore(documents_dir)
        _document_qa = DocumentQA(
            vector_store_manager=_vector_store_manager,
            llama_model_path=llama_model_path
        )
    
    return _document_qa

def get_document_qa() -> Optional[DocumentQA]:
    """Get the global DocumentQA instance."""
    return _document_qa