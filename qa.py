"""
Complete Document QA System
A single, simplified module that handles all document-based question answering.
"""

import os
import logging
import threading
import time
import json
import re
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

# Document processing
import pptx
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

# Llama integration (optional)
try:
    from llama_cpp import Llama
    LLAMA_AVAILABLE = True
except ImportError:
    LLAMA_AVAILABLE = False
    print("Info: llama-cpp-python not installed. AI will use rule-based responses.")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class ConversationMemory:
    """Simple conversation memory for context"""
    
    def __init__(self, max_history: int = 10):
        self.max_history = max_history
        self.conversations = {}
    
    def add_message(self, user_id: str, question: str, answer: str):
        if user_id not in self.conversations:
            self.conversations[user_id] = []
        
        self.conversations[user_id].append({
            'question': question,
            'answer': answer,
            'timestamp': datetime.now().isoformat()
        })
        
        # Keep only recent history
        if len(self.conversations[user_id]) > self.max_history:
            self.conversations[user_id] = self.conversations[user_id][-self.max_history:]
    
    def get_context(self, user_id: str, last_n: int = 3) -> str:
        if user_id not in self.conversations:
            return ""
        
        recent = self.conversations[user_id][-last_n:]
        context_parts = []
        
        for item in recent:
            context_parts.append(f"Human: {item['question']}")
            context_parts.append(f"Assistant: {item['answer']}")
        
        return "\n".join(context_parts)


class LlamaLLM:
    """Enhanced Llama model wrapper for Llama 3"""
    
    def __init__(self, model_path: str):
        if not LLAMA_AVAILABLE:
            raise ImportError("llama-cpp-python is required")
        
        self.model_path = model_path
        self.llm = None
        self._initialize_model()
    
    def _initialize_model(self):
        try:
            logger.info(f"Loading Llama model from {self.model_path}")
            self.llm = Llama(
                model_path=self.model_path,
                n_ctx=4096,           # Context window size
                n_threads=8,          # Number of CPU threads to use
                n_gpu_layers=35,      # Number of layers to offload to GPU (set to 0 for CPU only)
                verbose=False,
                temperature=0.2,      # Lower temperature for more deterministic outputs
                max_tokens=1024,      # Maximum tokens in response
                top_p=0.9,            # Top-p sampling
                top_k=40,             # Top-k sampling
                repeat_penalty=1.1    # Penalty for repeating tokens
            )
            logger.info("Llama model loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load Llama model: {str(e)}")
            raise
    
    def generate_response(self, prompt: str) -> str:
        if not self.llm:
            raise RuntimeError("Llama model not initialized")
        
        try:
            # For Llama 3, we'll generate a response with proper stop tokens
            response = self.llm(
                prompt,
                stop=["<|eot_id|>", "<|start_header_id|>user", "</s>"],  # Stop tokens for Llama 3
                echo=False           # Don't include the prompt in the response
            )
            
            # Extract and clean the response text
            text = response['choices'][0]['text'].strip()
            
            # Remove any trailing artifacts that might be in the response
            if "<|eot_id|>" in text:
                text = text.split("<|eot_id|>")[0].strip()
                
            return text
        except Exception as e:
            logger.error(f"Error generating response: {str(e)}")
            return "I apologize, but I encountered an error while generating a response."


class PPTXLoader:
    """Simple PPTX loader"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Dict[str, Any]]:
        docs = []
        try:
            prs = pptx.Presentation(self.file_path)
            filename = os.path.basename(self.file_path)
            
            for i, slide in enumerate(prs.slides):
                text_content = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_content += cell.text + " "
                            text_content += "\n"
                
                if text_content.strip():
                    doc = {
                        "page_content": text_content.strip(),
                        "metadata": {
                            "source": self.file_path,
                            "slide_number": i + 1,
                            "page": i + 1,
                            "filename": filename
                        }
                    }
                    docs.append(doc)
                    
            logger.info(f"Loaded {len(docs)} slides from {filename}")
            return docs
            
        except Exception as e:
            logger.error(f"Error processing PPTX {self.file_path}: {str(e)}")
            return []


class DocumentProcessor:
    """Document processor for all supported formats"""
    
    def __init__(self, documents_dir: str):
        self.documents_dir = documents_dir
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=100,
            length_function=len
        )

    def load_documents(self) -> List[Any]:
        all_docs = []
        supported_extensions = {'.pdf', '.pptx', '.ppt', '.txt'}
        
        if not os.path.exists(self.documents_dir):
            logger.warning(f"Documents directory not found: {self.documents_dir}")
            return []
        
        for root, _, files in os.walk(self.documents_dir):
            for file in files:
                file_path = os.path.join(root, file)
                filename = os.path.basename(file_path)
                
                try:
                    if file.lower().endswith('.pdf'):
                        docs = self._load_pdf(file_path, filename)
                        all_docs.extend(docs)
                        
                    elif file.lower().endswith(('.pptx', '.ppt')):
                        loader = PPTXLoader(file_path)
                        docs = loader.load()
                        all_docs.extend(docs)
                        
                    elif file.lower().endswith('.txt'):
                        docs = self._load_txt(file_path, filename)
                        all_docs.extend(docs)
                        
                except Exception as e:
                    logger.error(f"Error processing {file_path}: {str(e)}")
                    continue
        
        logger.info(f"Loaded {len(all_docs)} document sections")
        return all_docs

    def _load_pdf(self, file_path: str, filename: str) -> List[Any]:
        try:
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            
            for doc in docs:
                doc.metadata['filename'] = filename
                
            return docs
        except Exception as e:
            logger.error(f"Error loading PDF {filename}: {str(e)}")
            return []

    def _load_txt(self, file_path: str, filename: str) -> List[Any]:
        encodings = ['utf-8', 'utf-16', 'iso-8859-1', 'windows-1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                if content.strip():
                    return [{
                        "page_content": content.strip(),
                        "metadata": {
                            "source": file_path,
                            "page": 1,
                            "filename": filename
                        }
                    }]
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error loading TXT {filename}: {str(e)}")
                break
        return []

    def process_documents(self) -> List[Any]:
        docs = self.load_documents()
        
        if not docs:
            logger.warning("No documents loaded")
            return []
        
        logger.info(f"Processing {len(docs)} documents into chunks...")
        
        try:
            texts = []
            metadatas = []
            
            for doc in docs:
                if hasattr(doc, "page_content"):
                    content = doc.page_content
                    metadata = doc.metadata
                else:
                    content = doc["page_content"]
                    metadata = doc["metadata"]
                
                # Clean content
                content = re.sub(r'\s+', ' ', content).strip()
                if len(content) > 50:
                    texts.append(content)
                    metadatas.append(metadata)
            
            splits = self.text_splitter.create_documents(texts, metadatas=metadatas)
            logger.info(f"Created {len(splits)} document chunks")
            return splits
            
        except Exception as e:
            logger.error(f"Error processing documents: {str(e)}")
            return []


class VectorStore:
    """Vector store manager with background initialization"""
    
    def __init__(self, documents_dir: str, vector_db_path: str = "vector_db"):
        self.documents_dir = documents_dir
        self.vector_db_path = vector_db_path
        self.embeddings = None
        self.vector_store_instance = None
        self._initialization_lock = threading.Lock()
        self._is_initializing = False
        self._initialization_complete = False
        self._initialization_error = None
        
        self._start_background_initialization()

    def _start_background_initialization(self):
        def initialize():
            try:
                with self._initialization_lock:
                    if self._initialization_complete or self._is_initializing:
                        return
                    
                    self._is_initializing = True
                    logger.info("Starting vector store initialization...")
                    
                    self.embeddings = HuggingFaceEmbeddings(
                        model_name="all-MiniLM-L6-v2",
                        model_kwargs={'device': 'cpu'},
                        encode_kwargs={'normalize_embeddings': True}
                    )
                    
                    self.vector_store_instance = self._load_or_create_vector_store()
                    
                    self._initialization_complete = True
                    self._is_initializing = False
                    logger.info("Vector store initialization completed")
                    
            except Exception as e:
                self._initialization_error = str(e)
                self._is_initializing = False
                logger.error(f"Failed to initialize vector store: {str(e)}")
        
        thread = threading.Thread(target=initialize, daemon=True)
        thread.start()

    def _load_or_create_vector_store(self):
        if os.path.exists(self.vector_db_path):
            try:
                logger.info(f"Loading existing vector database from {self.vector_db_path}")
                return FAISS.load_local(
                    self.vector_db_path, 
                    self.embeddings, 
                    allow_dangerous_deserialization=True
                )
            except Exception as e:
                logger.warning(f"Failed to load existing vector store: {str(e)}")
        
        return self._create_new_vector_store()

    def _create_new_vector_store(self):
        processor = DocumentProcessor(self.documents_dir)
        document_chunks = processor.process_documents()
        
        if not document_chunks:
            logger.warning("No document chunks available")
            return None
        
        logger.info(f"Creating vector database from {len(document_chunks)} chunks...")
        
        try:
            vector_store = FAISS.from_documents(
                documents=document_chunks,
                embedding=self.embeddings
            )
            vector_store.save_local(self.vector_db_path)
            logger.info(f"Vector database saved to {self.vector_db_path}")
            return vector_store
        except Exception as e:
            logger.error(f"Failed to create vector database: {str(e)}")
            return None

    def get_vector_store(self, timeout: int = 30):
        start_time = time.time()
        
        while time.time() - start_time < timeout:
            if self._initialization_complete:
                return self.vector_store_instance
            elif self._initialization_error:
                raise Exception(f"Vector store failed: {self._initialization_error}")
            time.sleep(0.5)
        
        raise TimeoutError("Vector store initialization timed out")

    def is_ready(self) -> bool:
        return self._initialization_complete and self.vector_store_instance is not None


class DocumentQA:
    """Main Document QA system with conversational capabilities"""
    
    def __init__(self, documents_dir: str, llama_model_path: str = None):
        self.documents_dir = documents_dir
        self.vector_store_manager = VectorStore(documents_dir)
        self.llama_llm = None
        self.conversation_memory = ConversationMemory()
        
        # Initialize Llama if available
        if llama_model_path and LLAMA_AVAILABLE and os.path.exists(llama_model_path):
            try:
                self.llama_llm = LlamaLLM(llama_model_path)
                logger.info("Llama model initialized")
            except Exception as e:
                logger.error(f"Failed to initialize Llama: {str(e)}")
        
        # Greeting responses
        self.greetings = {
            'hi': "Hello! I'm your AI learning assistant. I can help you with course materials, explain concepts, or just chat about your studies. What would you like to know?",
            'hello': "Hi there! I'm here to help with your learning journey. You can ask me about course topics, request explanations, or get study tips. How can I assist you today?",
            'hey': "Hey! Great to see you here. I'm your AI tutor ready to help with anything you're studying. What's on your mind?",
            'good morning': "Good morning! Ready to learn something new today? I'm here to help with your course materials.",
            'good afternoon': "Good afternoon! How's your learning going today? I'm here to help with explanations or questions.",
            'good evening': "Good evening! Perfect time for some learning. What would you like to explore?"
        }

    def answer_question(self, question: str, user_id: str = None) -> Dict[str, Any]:
        """Main method to answer questions with conversational support"""
        
        if not question or not question.strip():
            return {
                "answer": "I'm here and ready to help! What would you like to know or discuss?",
                "sources": [],
                "conversation_type": "greeting"
            }
        
        question = question.strip()
        question_lower = question.lower()
        
        try:
            # Detect conversation type
            conversation_type = self._detect_conversation_type(question_lower)
            
            if conversation_type == "greeting":
                response = self._handle_greeting(question_lower)
                if user_id:
                    self.conversation_memory.add_message(user_id, question, response)
                return {
                    "answer": response,
                    "sources": [],
                    "conversation_type": "greeting"
                }
            
            # For other types, try document-based approach first
            response = self._handle_document_question(question, user_id)
            
            # If no good document answer and Llama available, enhance with general knowledge
            if (not response["sources"] and 
                self.llama_llm and 
                "couldn't find" in response["answer"].lower()):
                
                general_response = self._handle_general_question(question, user_id)
                if general_response["answer"] != response["answer"]:
                    response = {
                        "answer": general_response["answer"] + "\n\n*Note: This answer combines general knowledge with course materials.*",
                        "sources": response["sources"],
                        "conversation_type": "hybrid"
                    }
            
            if user_id:
                self.conversation_memory.add_message(user_id, question, response["answer"])
            
            return response
            
        except Exception as e:
            logger.error(f"Error in answer_question: {str(e)}")
            return {
                "answer": "I encountered an error, but I'm still here to help! Could you try rephrasing your question?",
                "sources": [],
                "conversation_type": "error"
            }

    def _detect_conversation_type(self, question_lower: str) -> str:
        # Check for greetings
        greeting_patterns = [
            r'\b(hi|hello|hey|hiya)\b',
            r'\bgood (morning|afternoon|evening|night)\b',
            r'\bhow are you\b',
            r'\bwhat\'?s up\b'
        ]
        
        if any(re.search(pattern, question_lower) for pattern in greeting_patterns):
            return "greeting"
        
        # Check for document-specific keywords
        document_keywords = [
            'explain', 'definition', 'what is', 'how does', 'tell me about',
            'course material', 'lesson', 'chapter', 'according to'
        ]
        
        if any(keyword in question_lower for keyword in document_keywords):
            return "document_based"
        
        return "general"

    def _handle_greeting(self, question_lower: str) -> str:
        for greeting, response in self.greetings.items():
            if greeting in question_lower:
                return response
        
        return "Hello! I'm your AI learning assistant. I'm here to help you understand course materials, answer questions, or discuss topics you're studying. What would you like to explore today?"

    def _handle_document_question(self, question: str, user_id: str = None) -> Dict[str, Any]:
        # Get vector store
        try:
            if not self.vector_store_manager.is_ready():
                if self.vector_store_manager._is_initializing:
                    return {
                        "answer": "I'm currently loading course materials. Please wait a moment and try again.",
                        "sources": [],
                        "conversation_type": "loading"
                    }
                else:
                    return {
                        "answer": "Course materials are not available. Please contact your administrator.",
                        "sources": [],
                        "conversation_type": "error"
                    }
            
            vector_store = self.vector_store_manager.get_vector_store(timeout=10)
            
            if not vector_store:
                return {
                    "answer": "I don't have access to course materials right now. Please ensure documents are uploaded.",
                    "sources": [],
                    "conversation_type": "error"
                }
            
            # Perform two queries to improve retrieval quality
            # First, retrieve based on semantic similarity
            retriever = vector_store.as_retriever(search_kwargs={"k": 6})
            docs = retriever.get_relevant_documents(question)
            
            # Then, perform a second retrieval with keywords from the question
            # This helps catch relevant documents that might be missed by semantic search
            keywords = self._extract_keywords(question)
            if keywords:
                keyword_query = " ".join(keywords)
                additional_docs = retriever.get_relevant_documents(keyword_query)
                
                # Combine the results while avoiding duplicates
                seen_contents = {doc.page_content for doc in docs}
                for doc in additional_docs:
                    if doc.page_content not in seen_contents:
                        docs.append(doc)
                        seen_contents.add(doc.page_content)
            
            if not docs:
                return {
                    "answer": "I couldn't find specific information about that in your course materials. Could you provide more context or try rephrasing your question?",
                    "sources": [],
                    "conversation_type": "document_based"
                }
            
            # Prepare context and sources
            context = self._prepare_context(docs)
            sources = self._extract_sources(docs)
            
            # Generate response
            if self.llama_llm:
                answer = self._generate_llama_response(question, context, user_id)
            else:
                answer = self._generate_rule_based_response(context, question)
            
            return {
                "answer": answer,
                "sources": sources,
                "conversation_type": "document_based"
            }
            
        except Exception as e:
            logger.error(f"Document question error: {str(e)}")
            return {
                "answer": "I encountered an error while searching course materials. Please try again.",
                "sources": [],
                "conversation_type": "error"
            }

    def _extract_keywords(self, question: str) -> List[str]:
        """Extract important keywords from a question for better retrieval"""
        import string
        
        # Common stop words to exclude
        stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 
            'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'being',
            'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could',
            'should', 'may', 'might', 'must', 'can', 'what', 'where', 'when',
            'why', 'how', 'who', 'which', 'this', 'that', 'these', 'those',
            'i', 'you', 'he', 'she', 'it', 'we', 'they', 'me', 'him', 'her',
            'us', 'them', 'my', 'your', 'his', 'her', 'its', 'our', 'their'
        }
        
        # Clean and split the question
        question_clean = question.lower().translate(str.maketrans('', '', string.punctuation))
        words = question_clean.split()
        
        # Filter out stop words and short words
        keywords = [word for word in words if word not in stop_words and len(word) > 2]
        
        return keywords[:5]  # Return top 5 keywords

    def _handle_general_question(self, question: str, user_id: str = None) -> Dict[str, Any]:
        if not self.llama_llm:
            return {
                "answer": "I'd be happy to help! Could you be more specific or ask about topics from your course materials?",
                "sources": [],
                "conversation_type": "general"
            }
        
        try:
            prompt = self._create_general_prompt(question, user_id)
            response = self.llama_llm.generate_response(prompt)
            
            return {
                "answer": response,
                "sources": [],
                "conversation_type": "general"
            }
        except Exception as e:
            logger.error(f"General question error: {str(e)}")
            return {
                "answer": "I encountered an error. Could you try rephrasing your question?",
                "sources": [],
                "conversation_type": "error"
            }

    def _prepare_context(self, docs, max_length: int = 3000) -> str:
        context_parts = []
        total_length = 0
        
        for doc in docs:
            content = doc.page_content if hasattr(doc, "page_content") else doc.get("page_content", "")
            if total_length + len(content) > max_length:
                break
            context_parts.append(content)
            total_length += len(content)
        
        return "\n\n".join(context_parts)

    def _extract_sources(self, docs) -> List[Dict[str, str]]:
        sources = []
        seen_sources = set()
        
        for doc in docs:
            metadata = doc.metadata if hasattr(doc, "metadata") else doc.get("metadata", {})
            filename = metadata.get("filename", os.path.basename(metadata.get("source", "Unknown")))
            page = metadata.get("page", metadata.get("slide_number", "Unknown"))
            
            source_key = f"{filename}_{page}"
            if source_key not in seen_sources:
                sources.append({
                    "file": filename,
                    "page": str(page)
                })
                seen_sources.add(source_key)
        
        return sources[:4]

    def _generate_llama_response(self, question: str, context: str, user_id: str = None) -> str:
        prompt = f"""<|begin_of_text|><|start_header_id|>system<|end_header_id|>

You are a helpful AI learning assistant for a learning platform. You provide clear, accurate, and educational responses based on course materials.

Guidelines:
- Use the course materials to provide accurate, helpful answers
- Be conversational and friendly in your tone
- Provide specific examples from the course materials when relevant
- If the materials don't fully answer the question, acknowledge this clearly
- Keep responses educational and informative, tailored to the learning context<|eot_id|>

<|start_header_id|>user<|end_header_id|>

I'm studying and need help understanding a topic from my course materials. Here are the relevant sections from my course materials:
{context}

My question is: {question}<|eot_id|>

<|start_header_id|>assistant<|end_header_id|>

"""
        
        try:
            return self.llama_llm.generate_response(prompt)
        except Exception as e:
            logger.error(f"Llama generation error: {str(e)}")
            return self._generate_rule_based_response(context, question)

    def _create_general_prompt(self, question: str, user_id: str = None) -> str:
        context = ""
        if user_id:
            context = self.conversation_memory.get_context(user_id, last_n=2)
        
        prompt = f"""<|begin_of_text|><|start_header_id|>system<|end_header_id|>

You are a friendly, helpful AI learning assistant for a learning platform. You provide conversational, engaging responses that are educational and supportive.

- Be helpful, friendly, and educational in your responses
- Explain concepts clearly using examples when possible
- Focus on being accurate and engaging
- Give practical advice when asked for help with studying or learning
- Always maintain a supportive, encouraging tone for students

{f"Previous conversation:\n{context}" if context else ""}<|eot_id|>

<|start_header_id|>user<|end_header_id|>

{question}<|eot_id|>

<|start_header_id|>assistant<|end_header_id|>

"""
        return prompt

    def _generate_rule_based_response(self, context: str, question: str) -> str:
        # Simple rule-based fallback
        sentences = [s.strip() for s in context.split('.') if s.strip() and len(s.strip()) > 20]
        question_words = set(question.lower().split())
        
        relevant_sentences = []
        for sentence in sentences[:10]:
            sentence_words = set(sentence.lower().split())
            overlap = len(question_words.intersection(sentence_words))
            if overlap > 0:
                relevant_sentences.append((sentence, overlap))
        
        relevant_sentences.sort(key=lambda x: x[1], reverse=True)
        top_sentences = [s[0] for s in relevant_sentences[:3]]
        
        if top_sentences:
            return "Based on your course materials:\n\n" + ". ".join(top_sentences) + "."
        else:
            return "I found some information in the course materials, but couldn't extract a specific answer. Could you be more specific about what you'd like to know?"

    def get_status(self) -> Dict[str, Any]:
        """Get system status"""
        document_count = 0
        if os.path.exists(self.documents_dir):
            for root, dirs, files in os.walk(self.documents_dir):
                for file in files:
                    if file.lower().endswith(('.pdf', '.pptx', '.ppt', '.txt')):
                        document_count += 1
        
        return {
            "ready": self.vector_store_manager.is_ready(),
            "initializing": self.vector_store_manager._is_initializing,
            "error": self.vector_store_manager._initialization_error,
            "document_count": document_count,
            "llama_available": self.llama_llm is not None
        }


# Global instance management
_qa_instance = None

def initialize_qa(documents_dir: str = "documents", llama_model_path: str = None) -> DocumentQA:
    """Initialize the QA system"""
    global _qa_instance
    
    if _qa_instance is None:
        logger.info(f"Initializing QA system with documents from: {documents_dir}")
        _qa_instance = DocumentQA(documents_dir, llama_model_path)
    
    return _qa_instance

def get_qa_system() -> Optional[DocumentQA]:
    """Get the global QA instance"""
    return _qa_instance

def ask_question(question: str, documents_dir: str = "documents", llama_model_path: str = None, user_id: str = None) -> Dict[str, Any]:
    """Simple function to ask a question"""
    qa_system = initialize_qa(documents_dir, llama_model_path)
    return qa_system.answer_question(question, user_id)

def get_system_status(documents_dir: str = "documents", llama_model_path: str = None) -> Dict[str, Any]:
    """Get system status"""
    qa_system = initialize_qa(documents_dir, llama_model_path)
    return qa_system.get_status()

# Example usage
if __name__ == "__main__":
    # Initialize QA system
    qa = DocumentQA("documents")
    
    # Check status
    status = qa.get_status()
    print("QA System Status:", status)
    
    if status["ready"]:
        # Ask questions
        response = qa.answer_question("Hello!")
        print("\nGreeting Response:", response)
        
        response = qa.answer_question("What is artificial intelligence?")
        print("\nAI Question Response:", response)
    else:
        print("QA system is not ready. Add documents to the 'documents' directory.")
